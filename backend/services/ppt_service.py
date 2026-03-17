"""
PPT 텍스트 추출 서비스

- 텍스트 상자 모드: #숫자 또는 #숫자-숫자 마커가 포함된 텍스트 상자 추출
- 메모 모드: 슬라이드 메모 전체 추출
"""
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

from config import PATTERNS
from models.schemas import SlideInfo, TextBoxCandidate


# 슬라이드 마커 패턴 (#1, #3-1 등)
MARKER_RE = re.compile(PATTERNS["slide_marker"]["pattern"])


def extract_slides(pptx_path: str | Path, extraction_location: str) -> list[SlideInfo]:
    """
    PPT 파일에서 슬라이드 정보를 추출합니다.

    Args:
        pptx_path: .pptx 파일 경로
        extraction_location: "textbox" 또는 "notes"

    Returns:
        SlideInfo 리스트 (슬라이드별 후보 텍스트 상자 또는 자동 선택된 텍스트)
    """
    prs = Presentation(str(pptx_path))
    slides_info: list[SlideInfo] = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_number = slide_idx + 1

        if extraction_location == "notes":
            # ── 메모 추출 모드 ──────────────────────────────────────────
            notes_text = _extract_notes(slide)
            # 메모는 단일 텍스트이므로 자동 선택
            candidate = TextBoxCandidate(
                id=0,
                marker="notes",
                preview=notes_text[:60] + ("..." if len(notes_text) > 60 else ""),
                full_text=notes_text,
            )
            slides_info.append(SlideInfo(
                slide_number=slide_number,
                candidates=[candidate],
                auto_selected=True,
                selected_text=notes_text,
            ))

        else:
            # ── 텍스트 상자 추출 모드 ───────────────────────────────────
            candidates = _extract_textbox_candidates(slide)

            if len(candidates) == 0:
                # 마커가 있는 텍스트 상자가 없으면 빈 슬라이드로 처리
                slides_info.append(SlideInfo(
                    slide_number=slide_number,
                    candidates=[],
                    auto_selected=True,
                    selected_text="",
                ))
            elif len(candidates) == 1:
                # 후보 1개 → 자동 선택
                slides_info.append(SlideInfo(
                    slide_number=slide_number,
                    candidates=candidates,
                    auto_selected=True,
                    selected_text=candidates[0].full_text,
                ))
            else:
                # 후보 여러 개 → 사용자 선택 필요
                slides_info.append(SlideInfo(
                    slide_number=slide_number,
                    candidates=candidates,
                    auto_selected=False,
                    selected_text=None,
                ))

    return slides_info


def _extract_textbox_candidates(slide) -> list[TextBoxCandidate]:
    """슬라이드에서 #마커가 포함된 텍스트 상자를 추출합니다."""
    candidates = []
    textbox_idx = 0

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        # 텍스트 상자의 전체 텍스트 추출 (단락 결합)
        full_text = _extract_shape_text(shape)

        # #숫자 또는 #숫자-숫자 패턴 검색
        match = MARKER_RE.search(full_text)
        if match:
            marker = match.group(0)  # 예: "#1", "#3-1"
            preview = full_text[:80].replace("\n", " ")
            candidates.append(TextBoxCandidate(
                id=textbox_idx,
                marker=marker,
                preview=preview,
                full_text=full_text,
            ))

        textbox_idx += 1

    return candidates


def _extract_shape_text(shape) -> str:
    """텍스트 프레임에서 모든 단락의 텍스트를 추출합니다."""
    lines = []
    for para in shape.text_frame.paragraphs:
        # 단락 내 모든 run의 텍스트를 결합
        line = "".join(run.text for run in para.runs)
        if line.strip():
            lines.append(line)
    return "\n".join(lines)


def _extract_notes(slide) -> str:
    """슬라이드 메모를 추출합니다."""
    if not slide.has_notes_slide:
        return ""
    notes_slide = slide.notes_slide
    if not notes_slide.notes_text_frame:
        return ""
    lines = []
    for para in notes_slide.notes_text_frame.paragraphs:
        line = "".join(run.text for run in para.runs)
        if line.strip():
            lines.append(line)
    return "\n".join(lines)
