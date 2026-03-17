"""
테스트용 샘플 .pptx 파일 생성 스크립트

실행: python create_sample_pptx.py
→ sample_04차시.pptx 생성
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def create_sample():
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # 빈 레이아웃

    slides_data = [
        {
            "marker_text": "#1 AI 리스크 관리 개요",
            "body_text": "나레이션: AI(Artificial Intelligence) 시스템의 리스크 관리는 2026년 현재 ChatGPT와 같은 LLM 도입이 늘면서 중요해졌습니다. (Ver.2) 기준으로 설명합니다.",
            "notes": "오늘 강의의 첫 번째 슬라이드입니다. AI 리스크 개요를 다룹니다.",
        },
        {
            "marker_text": "#2 AI 리스크의 유형",
            "body_text": "나레이션: AI 리스크 유형에는 Bias, Hallucination, Privacy 침해 등이 있습니다. ISO/IEC 42001 표준을 참고하십시오.",
            "notes": "두 번째 슬라이드: 리스크 유형 설명.",
        },
        {
            "marker_text": "#3 관리 프레임워크",
            "body_text": "나레이션: NIST AI RMF(Risk Management Framework)와 EU AI Act를 기반으로 합니다. 총 4단계(Govern, Map, Measure, Manage)로 구성됩니다.",
            "notes": "세 번째 슬라이드: 관리 프레임워크 소개.",
        },
        {
            "marker_text": "#4 실습: ChatGPT 활용",
            "body_text": "나레이션: ChatGPT API(v4.0)를 활용하여 리스크 평가 보고서를 자동 생성합니다. Python 3.11 환경에서 실습합니다.",
            "notes": "네 번째 슬라이드: ChatGPT 실습.",
        },
        {
            "marker_text": "#5 정리 및 Q&A",
            "body_text": "나레이션: 오늘 강의를 정리합니다. AI 리스크 관리의 핵심은 Transparency와 Accountability입니다. Q&A 시간을 갖겠습니다.",
            "notes": "다섯 번째 슬라이드: 정리.",
        },
    ]

    for data in slides_data:
        slide = prs.slides.add_slide(slide_layout)

        # 마커가 포함된 텍스트 상자 (본문)
        tx_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(1)
        )
        tf = tx_box.text_frame
        tf.text = data["marker_text"]
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.size = Pt(18)

        # 본문 텍스트 상자
        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8), Inches(9), Inches(4)
        )
        body_tf = body_box.text_frame
        body_tf.text = data["body_text"]
        body_tf.paragraphs[0].runs[0].font.size = Pt(14)

        # 메모 추가
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = data["notes"]

    prs.save("sample_04차시.pptx")
    print("✅ sample_04차시.pptx 생성 완료")
    print("   슬라이드 수:", len(slides_data))
    print("   접두어 '나레이션:' 포함")
    print("   비한글 항목: AI, ChatGPT, LLM, ISO/IEC 42001, NIST, ...")


if __name__ == "__main__":
    create_sample()
